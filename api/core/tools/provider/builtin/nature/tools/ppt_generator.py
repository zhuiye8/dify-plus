"""
PPT生成工具实现
"""
from typing import Any, Union, Optional
import os
import json
import base64
from io import BytesIO
import tempfile

from pptx import Presentation

from core.tools.entities.tool_entities import ToolInvokeMessage
from core.tools.tool.builtin_tool import BuiltinTool
from core.file.file_manager import download

class PPTGeneratorTool(BuiltinTool):
    """
    PPT生成工具类
    
    根据大纲文本和PPT模板生成演示文稿
    """

    def _get_file_content(self, file_obj: Any) -> bytes:
        """
        从文件对象中获取文件内容
        """
        try:
            # 如果是 core.file.models.File 类型
            if hasattr(file_obj, 'path') and os.path.exists(file_obj.path):
                with open(file_obj.path, 'rb') as f:
                    return f.read()
            
            # 尝试使用 download 函数
            file_content = download(file_obj)
            if file_content:
                return file_content
                
            raise ValueError(f"不支持的文件对象类型: {type(file_obj)}")
            
        except Exception as e:
            raise ValueError(f"读取文件内容失败: {str(e)}")

    def _invoke(
        self,
        user_id: str,
        tool_parameters: dict[str, Any],
    ) -> Union[ToolInvokeMessage, list[ToolInvokeMessage]]:
        """
        执行PPT生成逻辑
        """
        outline = tool_parameters.get("outline")
        template_file = tool_parameters.get("template")

        # 校验输入
        if not outline:
            raise Exception("PPT大纲未提供")

        try:
            # 如果提供了模板，读取模板文件
            if template_file:
                template_content = self._get_file_content(template_file)
                
                # 创建临时文件保存模板
                with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
                    temp_file.write(template_content)
                    template_path = temp_file.name
            else:
                # 使用默认空白模板
                prs = Presentation()
                with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
                    prs.save(temp_file.name)
                    template_path = temp_file.name
            
            # 生成PPT
            result_ppt_io = self.generate_ppt_from_outline(outline, template_path)
            
            # 删除临时文件
            if os.path.exists(template_path):
                os.unlink(template_path)
            
            # 读取生成的PPT内容
            result_ppt_io.seek(0)
            ppt_content = result_ppt_io.read()
            
            # 生成文件名
            output_filename = "generated_presentation.pptx"
            if template_file and hasattr(template_file, 'name'):
                base_name = os.path.splitext(template_file.name)[0]
                output_filename = f"{base_name}_generated.pptx"
                
            return [
                self.create_text_message("PPT生成成功"),
                self.create_blob_message(
                    blob=ppt_content,
                    meta={
                        "mime_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        "name": output_filename
                    },
                    save_as=self.VariableKey.PPT
                )
            ]
            
        except Exception as e:
            return self.create_text_message(f"PPT生成失败: {str(e)}")

    def generate_ppt_from_outline(self, outline_text, template_path):
        """
        根据大纲文本和PPT模板生成新的PPT文件
        
        参数:
        outline_text (str): PPT大纲文本，格式如"# 标题\n## 子标题\n- 要点1\n- 要点2"
        template_path (str): PPT模板文件路径
        
        返回:
        BytesIO: 包含生成的PPT文件的二进制数据
        """
        # 加载模板
        try:
            prs = Presentation(template_path)
        except Exception as e:
            raise Exception(f"无法加载PPT模板: {str(e)}")
        
        # 解析大纲文本
        slides_content = self.parse_outline(outline_text)
        
        # 为每个解析出的内容创建幻灯片
        for slide_content in slides_content:
            # 获取适当的布局
            layout = self.get_suitable_layout(prs, slide_content)
            
            # 创建新幻灯片
            slide = prs.slides.add_slide(layout)
            
            # 填充幻灯片内容
            self.fill_slide_content(slide, slide_content)
        
        # 保存到内存中的文件对象
        ppt_io = BytesIO()
        prs.save(ppt_io)
        ppt_io.seek(0)
        
        return ppt_io

    def parse_outline(self, outline_text):
        """
        解析大纲文本，提取幻灯片内容
        
        参数:
        outline_text (str): 大纲文本
        
        返回:
        list: 包含每张幻灯片内容的字典列表
        """
        slides = []
        current_slide = None
        
        # 按行分割大纲
        lines = outline_text.strip().split('\n')
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            # 主标题 (# 开头)
            if line.startswith('# '):
                if current_slide:
                    slides.append(current_slide)
                current_slide = {
                    'title': line[2:].strip(),
                    'subtitle': '',
                    'content': []
                }
            
            # 副标题 (## 开头)
            elif line.startswith('## '):
                if current_slide:
                    current_slide['subtitle'] = line[3:].strip()
                else:
                    # 如果没有主标题就遇到副标题，创建一个新幻灯片
                    current_slide = {
                        'title': '',
                        'subtitle': line[3:].strip(),
                        'content': []
                    }
            
            # 内容项 (- 或* 开头)
            elif line.startswith('- ') or line.startswith('* '):
                if current_slide:
                    current_slide['content'].append(line[2:].strip())
                else:
                    # 如果没有标题就遇到内容，创建一个新幻灯片
                    current_slide = {
                        'title': '',
                        'subtitle': '',
                        'content': [line[2:].strip()]
                    }
            
            # 普通文本
            else:
                if current_slide:
                    # 如果前面有内容项，添加为子项
                    if current_slide['content']:
                        # 检查是否应该作为上一个要点的子内容
                        if line.startswith('  '):
                            last_item = current_slide['content'][-1]
                            current_slide['content'][-1] = f"{last_item}\n{line.strip()}"
                        else:
                            current_slide['content'].append(line)
                    else:
                        current_slide['content'].append(line)
                else:
                    # 如果没有标题就遇到普通文本，创建一个新幻灯片
                    current_slide = {
                        'title': '',
                        'subtitle': '',
                        'content': [line]
                    }
        
        # 添加最后一张幻灯片
        if current_slide:
            slides.append(current_slide)
        
        return slides

    def get_suitable_layout(self, prs, slide_content):
        """
        根据幻灯片内容选择合适的幻灯片布局
        
        参数:
        prs (Presentation): PPT演示文稿对象
        slide_content (dict): 幻灯片内容
        
        返回:
        SlideLayout: 幻灯片布局对象
        """
        # 获取可用的布局
        layouts = prs.slide_layouts
        
        # 根据内容选择布局
        if slide_content['title'] and slide_content['subtitle'] and slide_content['content']:
            # 标题、副标题和内容 - 寻找标题和内容布局
            for i, layout in enumerate(layouts):
                if "title and content" in layout.name.lower():
                    return layout
            
            # 如果找不到，使用标题布局
            for i, layout in enumerate(layouts):
                if "title" in layout.name.lower():
                    return layout
        
        elif slide_content['title'] and not slide_content['content']:
            # 只有标题 - 寻找标题布局
            for i, layout in enumerate(layouts):
                if "title" in layout.name.lower() and "content" not in layout.name.lower():
                    return layout
        
        elif slide_content['content'] and not slide_content['title']:
            # 只有内容 - 寻找内容布局
            for i, layout in enumerate(layouts):
                if "content" in layout.name.lower() and "title" not in layout.name.lower():
                    return layout
        
        # 默认使用第一个布局（通常是标题布局）
        return layouts[0]

    def fill_slide_content(self, slide, slide_content):
        """
        填充幻灯片内容
        
        参数:
        slide (Slide): 幻灯片对象
        slide_content (dict): 幻灯片内容
        """
        # 遍历占位符
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
                
            # 根据占位符类型填充内容
            if shape.placeholder_format.type == 1:  # 标题
                shape.text = slide_content['title']
            elif shape.placeholder_format.type == 2:  # 内容
                # 检查是否有内容项
                if slide_content['content']:
                    # 创建项目符号列表
                    text_frame = shape.text_frame
                    text_frame.clear()
                    
                    for i, item in enumerate(slide_content['content']):
                        p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
                        p.text = item
                        p.level = 0
            elif shape.placeholder_format.type == 3:  # 副标题
                shape.text = slide_content['subtitle']
            
            # 如果没有识别出占位符类型，尝试基于名称进行匹配
            elif "title" in str(shape.name).lower():
                shape.text = slide_content['title']
            elif "subtitle" in str(shape.name).lower() or "sub-title" in str(shape.name).lower():
                shape.text = slide_content['subtitle']
            elif "content" in str(shape.name).lower() or "body" in str(shape.name).lower():
                # 检查是否有内容项
                if slide_content['content']:
                    # 创建项目符号列表
                    text_frame = shape.text_frame
                    text_frame.clear()
                    
                    for i, item in enumerate(slide_content['content']):
                        p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
                        p.text = item
                        p.level = 0

    def validate_credentials(self, credentials: dict[str, Any], tool_parameters: dict[str, Any]) -> None:
        """
        验证凭证
        
        Args:
            credentials: 凭证字典
            tool_parameters: 工具参数
        """
        # 不需要验证凭证
        pass 