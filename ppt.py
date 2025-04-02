import os
from pptx import Presentation
import re
import tempfile
import base64
from io import BytesIO

def generate_ppt_from_outline(outline_text, template_path):
    """
    根据大纲文本和PPT模板生成新的PPT文件
    
    参数:
    outline_text (str): PPT大纲文本，格式如"# 标题\n## 子标题\n- 要点1\n- 要点2"
    template_path (str): PPT模板文件路径(.potx或.pptx)
    
    返回:
    BytesIO: 包含生成的PPT文件的二进制数据
    """
    # 加载模板
    try:
        prs = Presentation(template_path)
    except Exception as e:
        raise Exception(f"无法加载PPT模板: {str(e)}")
    
    # 解析大纲文本
    slides_content = parse_outline(outline_text)
    
    # 为每个解析出的内容创建幻灯片
    for slide_content in slides_content:
        # 获取适当的布局
        layout = get_suitable_layout(prs, slide_content)
        
        # 创建新幻灯片
        slide = prs.slides.add_slide(layout)
        
        # 填充幻灯片内容
        fill_slide_content(slide, slide_content)
    
    # 保存到内存中的文件对象
    ppt_io = BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    
    return ppt_io

def parse_outline(outline_text):
    """
    解析大纲文本，提取幻灯片内容
    
    参数:
    outline_text (str): 大纲文本
    
    返回:
    list: 包含每张幻灯片内容的字典列表
    """
    slides = []
    current_slide = None
    
    # 按行分割大纲
    lines = outline_text.strip().split('\n')
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
            
        # 主标题 (# 开头)
        if line.startswith('# '):
            if current_slide:
                slides.append(current_slide)
            current_slide = {
                'title': line[2:].strip(),
                'subtitle': '',
                'content': []
            }
        
        # 副标题 (## 开头)
        elif line.startswith('## '):
            if current_slide:
                current_slide['subtitle'] = line[3:].strip()
            else:
                # 如果没有主标题就遇到副标题，创建一个新幻灯片
                current_slide = {
                    'title': '',
                    'subtitle': line[3:].strip(),
                    'content': []
                }
        
        # 内容项 (- 或* 开头)
        elif line.startswith('- ') or line.startswith('* '):
            if current_slide:
                current_slide['content'].append(line[2:].strip())
            else:
                # 如果没有标题就遇到内容，创建一个新幻灯片
                current_slide = {
                    'title': '',
                    'subtitle': '',
                    'content': [line[2:].strip()]
                }
        
        # 普通文本
        else:
            if current_slide:
                # 如果前面有内容项，添加为子项
                if current_slide['content']:
                    # 检查是否应该作为上一个要点的子内容
                    if line.startswith('  '):
                        last_item = current_slide['content'][-1]
                        current_slide['content'][-1] = f"{last_item}\n{line.strip()}"
                    else:
                        current_slide['content'].append(line)
                else:
                    current_slide['content'].append(line)
            else:
                # 如果没有标题就遇到普通文本，创建一个新幻灯片
                current_slide = {
                    'title': '',
                    'subtitle': '',
                    'content': [line]
                }
    
    # 添加最后一张幻灯片
    if current_slide:
        slides.append(current_slide)
    
    return slides

def get_suitable_layout(prs, slide_content):
    """
    根据幻灯片内容选择合适的幻灯片布局
    
    参数:
    prs (Presentation): PPT演示文稿对象
    slide_content (dict): 幻灯片内容
    
    返回:
    SlideLayout: 幻灯片布局对象
    """
    # 获取可用的布局
    layouts = prs.slide_layouts
    
    # 根据内容选择布局
    if slide_content['title'] and slide_content['subtitle'] and slide_content['content']:
        # 标题、副标题和内容 - 寻找标题和内容布局
        for i, layout in enumerate(layouts):
            if "title and content" in layout.name.lower():
                return layout
        
        # 如果找不到，使用标题布局
        for i, layout in enumerate(layouts):
            if "title" in layout.name.lower():
                return layout
    
    elif slide_content['title'] and not slide_content['content']:
        # 只有标题 - 寻找标题布局
        for i, layout in enumerate(layouts):
            if "title" in layout.name.lower() and "content" not in layout.name.lower():
                return layout
    
    elif slide_content['content'] and not slide_content['title']:
        # 只有内容 - 寻找内容布局
        for i, layout in enumerate(layouts):
            if "content" in layout.name.lower() and "title" not in layout.name.lower():
                return layout
    
    # 默认使用第一个布局（通常是标题布局）
    return layouts[0]

def fill_slide_content(slide, slide_content):
    """
    填充幻灯片内容
    
    参数:
    slide (Slide): 幻灯片对象
    slide_content (dict): 幻灯片内容
    """
    # 遍历占位符
    for shape in slide.shapes:
        if not hasattr(shape, "text"):
            continue
            
        # 根据占位符类型填充内容
        if hasattr(shape, "placeholder_format") and hasattr(shape.placeholder_format, "type"):
            if shape.placeholder_format.type == 1:  # 标题
                shape.text = slide_content['title']
            elif shape.placeholder_format.type == 2:  # 内容
                # 检查是否有内容项
                if slide_content['content']:
                    # 创建项目符号列表
                    text_frame = shape.text_frame
                    text_frame.clear()
                    
                    for i, item in enumerate(slide_content['content']):
                        p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
                        p.text = item
                        p.level = 0
            elif shape.placeholder_format.type == 3:  # 副标题
                shape.text = slide_content['subtitle']
        
        # 如果没有识别出占位符类型，尝试基于名称进行匹配
        elif hasattr(shape, "name"):
            if "title" in str(shape.name).lower():
                shape.text = slide_content['title']
            elif "subtitle" in str(shape.name).lower() or "sub-title" in str(shape.name).lower():
                shape.text = slide_content['subtitle']
            elif "content" in str(shape.name).lower() or "body" in str(shape.name).lower():
                # 检查是否有内容项
                if slide_content['content']:
                    # 创建项目符号列表
                    text_frame = shape.text_frame
                    text_frame.clear()
                    
                    for i, item in enumerate(slide_content['content']):
                        p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
                        p.text = item
                        p.level = 0

def parse_base64_file(base64_string):
    """
    解析Base64编码的文件
    
    参数:
    base64_string (str): Base64编码的文件字符串
    
    返回:
    BytesIO: 解码后的文件对象
    """
    file_data = base64.b64decode(base64_string)
    return BytesIO(file_data)

def file_to_base64(file_io):
    """
    将文件对象转换为Base64编码
    
    参数:
    file_io (BytesIO): 文件对象
    
    返回:
    str: Base64编码的字符串
    """
    return base64.b64encode(file_io.getvalue()).decode('utf-8')

# 用于Dify工具集成的函数
def generate_ppt(outline, template_base64=None):
    """
    生成PPT的主函数，可以集成到Dify工具中
    
    参数:
    outline (str): PPT大纲文本
    template_base64 (str, optional): Base64编码的PPT模板文件(.potx或.pptx)
    
    返回:
    dict: 包含生成的PPT文件的Base64编码和文件名
    """
    try:
        # 如果提供了模板，解析Base64编码
        if template_base64:
            template_io = parse_base64_file(template_base64)
            
            # 创建临时文件保存模板
            # 检测文件格式并使用正确的后缀
            file_content = template_io.getvalue()
            file_suffix = '.potx' if file_content[0:4] == b'\xD0\xCF\x11\xE0' else '.pptx'
            
            with tempfile.NamedTemporaryFile(suffix=file_suffix, delete=False) as temp_file:
                temp_file.write(template_io.getvalue())
                template_path = temp_file.name
        else:
            # 使用默认空白模板
            from pptx import Presentation
            prs = Presentation()
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
                prs.save(temp_file.name)
                template_path = temp_file.name
        
        # 生成PPT
        result_ppt_io = generate_ppt_from_outline(outline, template_path)
        
        # 删除临时文件
        if os.path.exists(template_path):
            os.unlink(template_path)
        
        # 将结果转换为Base64
        result_base64 = file_to_base64(result_ppt_io)
        
        return {
            "success": True,
            "file_base64": result_base64,
            "filename": "generated_presentation.pptx"
        }
        
    except Exception as e:
        return {
            "success": False,
            "error": str(e)
        }

# 示例用法
if __name__ == "__main__":
    # 示例大纲
    sample_outline = """# 项目报告
## 2025年第一季度
- 项目概述
- 主要成就
- 遇到的挑战
- 下一步计划

# 项目概述
- 开始日期：2025年1月15日
- 结束日期：2025年3月31日
- 团队成员：5人
- 预算：100,000元

# 主要成就
- 完成了所有计划的功能开发
- 提前一周交付了MVP版本
- 获得了客户的积极反馈
- 团队协作效率提高30%

# 遇到的挑战
- 技术难点突破花费了额外时间
- 需求变更导致部分重构
- 资源短缺问题

# 下一步计划
- 发布1.0正式版
- 开始用户培训
- 收集反馈并优化
- 准备下一阶段计划
"""
    
    # 不使用模板，使用默认空白模板
    result = generate_ppt(sample_outline)
    
    if result["success"]:
        print(f"PPT已生成，文件名：{result['filename']}")
        # 可以将Base64编码保存为文件
        with open(result["filename"], "wb") as f:
            f.write(base64.b64decode(result["file_base64"]))
    else:
        print(f"PPT生成失败：{result['error']}")